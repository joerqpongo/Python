from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Crear presentación
prs = Presentation()

# Función para agregar diapositivas
def add_slide(layout_num, title_text, subtitle_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_num])
    title = slide.shapes.title
    title.text = title_text
    if subtitle_text:
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text

# Diapositiva 1: Título
add_slide(0, "Python en la Administración: Un Potencial Ilimitado")

# Diapositiva 2: ¿Qué es Python?
add_slide(1, "¿Qué es Python?", "Un lenguaje de programación versátil, fácil de aprender y potente.")

# Diapositiva 3: Ventajas a corto plazo
add_slide(1, "Ventajas a Corto Plazo")
bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = bullet_slide.shapes
body_shape = shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "Automatización de tareas repetitivas"
p = tf.add_paragraph()
p.text = "Análisis de datos rápido y eficiente"
p.level = 1
p = tf.add_paragraph()
p.text = "Creación de informes personalizados"
p.level = 1
p = tf.add_paragraph()
p.text = "Mejora en la toma de decisiones"
p.level = 1


# Diapositiva 4: Ventajas a largo plazo
add_slide(1, "Ventajas a Largo Plazo")
bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = bullet_slide.shapes
body_shape = shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "Desarrollo de herramientas internas a medida"
p = tf.add_paragraph()
p.text = "Optimización de procesos"
p.level = 1
p = tf.add_paragraph()
p.text = "Mayor eficiencia y productividad"
p.level = 1
p = tf.add_paragraph()
p.text = "Reducción de costes"
p.level = 1

# Diapositiva 5: Conclusión
add_slide(1, "Conclusión", "Python es una inversión estratégica para el departamento de administración.")

# Guardar la presentación
prs.save("C:/JL/presentación_python.pptx")
