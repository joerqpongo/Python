import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Crear la presentación
prs = Presentation()

# Personalizar el estilo de la presentación
def set_slide_style(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)  # Color de fondo gris claro

# Función para añadir diapositiva con título y contenido en bullet points
def add_slide_with_bullets(title, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_style(slide)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)  # Tamaño de fuente del título
    title_shape.text_frame.paragraphs[0].font.bold = True  # Negrita en el título
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Centrar el título

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Limpiar el marcador de posición

    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1  # Nivel de viñeta (1 para principal, 2 para secundaria, etc.)
        p.font.size = Pt(18)  # Tamaño de fuente de las viñetas

# Diapositiva de portada (con imagen de fondo)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Diseño de diapositiva en blanco
set_slide_style(slide)

# Insertar imagen de fondo (reemplaza 'ruta/a/tu/imagen.jpg' por la ruta real)
img_path = 'C:/JL/imagen presentación.jpeg'  
left = top = Inches(0)
pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)

# Agregar título y subtítulo sobre la imagen
title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
title.text = "Informe sobre el Sistema de Fichajes"
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # Color blanco

subtitle = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
subtitle.text = "Descripción y Configuraciones"
subtitle.text_frame.paragraphs[0].font.size = Pt(24)
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Resto de las diapositivas (utilizando la función add_slide_with_bullets)
# ... (el resto de tu código para agregar diapositivas)

# Guardar la presentación
pptx_file = os.path.join("C:/JL", "presentación.pptx")
prs.save(pptx_file)
